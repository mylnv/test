from flask import Flask, request, jsonify, render_template, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import os
from pptx import Presentation

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file:
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)

        # Convert PPTX to text
        try:
            text_file = convert_pptx_to_txt(file_path)
            return jsonify({'message': 'File successfully processed', 'download_url': f'/download/{text_file}'}), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    path = os.path.join(OUTPUT_FOLDER, filename)
    if os.path.exists(path):
        return send_file(path, as_attachment=True)
    else:
        return jsonify({'error': 'File not found'}), 404


def convert_pptx_to_txt(file_path):
    prs = Presentation(file_path)
    text_file = f"{os.path.basename(file_path).rsplit('.', 1)[0]}.txt"
    output_path = os.path.join(OUTPUT_FOLDER, text_file)

    with open(output_path, 'w', encoding='utf-8') as f:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        f.write(paragraph.text + '\n')
    
    return text_file

if __name__ == "__main__":
    app.run(debug=True)